import streamlit as st
import os
import re
import time
import ollama
from pptx import Presentation
from io import BytesIO

# --- CONFIGURATION ---
# Ensure you use a model you have locally. 
MODEL_NAME = "qwen3:30b" 
SEPARATOR = "@@@_START_SLIDE_CONTENT_@@@"

# --- CSS & AESTHETICS ---
def inject_derek_central_css():
    st.markdown("""
    <style>
    /* 1. RESET & FONTS */
    @import url('https://fonts.googleapis.com/css2?family=Satoshi:wght@400;500;700&display=swap');
    
    html, body, [class*="css"] {
        font-family: 'Satoshi', sans-serif;
    }

    /* 2. BACKGROUND & BOKEH */
    .stApp {
        background-color: #050505;
    }
    .stApp::before {
        content: ''; position: fixed; top: 0; left: 0; right: 0; bottom: 0; z-index: 0;
        background: radial-gradient(circle at 15% 50%, rgba(120, 0, 255, 0.1), transparent 25%),
                    radial-gradient(circle at 45% 80%, rgba(210, 0, 255, 0.1), transparent 25%),
                    radial-gradient(circle at 85% 30%, rgba(255, 0, 150, 0.1), transparent 25%);
        pointer-events: none;
    }

    /* 3. AURORA HEADER ANIMATION */
    @keyframes aurora-flow {
        0% { background-position: 0% 0%; }
        50% { background-position: 100% 100%; }
        100% { background-position: 0% 0%; }
    }
    
    .aurora-text {
        background: linear-gradient(300deg, #ffffff, #a5b4fc, #667eea, #764ba2, #f093fb, #667eea);
        background-size: 300% 300%;
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        animation: aurora-flow 6s ease infinite;
        font-weight: 700;
        font-size: 3.5rem;
        margin: 0;
        padding-bottom: 10px;
    }

    /* 4. BUTTON OVERRIDE */
    div.stButton > button {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%) !important;
        color: white !important;
        border: none !important;
        border-radius: 12px !important;
        padding: 0.75rem 1.5rem !important;
        font-weight: 600 !important;
        box-shadow: 0 4px 15px rgba(102, 126, 234, 0.4) !important;
        transition: transform 0.2s, box-shadow 0.2s !important;
    }
    div.stButton > button:hover {
        transform: scale(1.02);
        box-shadow: 0 6px 20px rgba(102, 126, 234, 0.6) !important;
    }

    /* 5. INPUT FIELDS & TEXTAREAS */
    [data-testid="stTextInput"] > div > div, 
    [data-testid="stTextArea"] > div > div {
        background-color: rgba(255, 255, 255, 0.03) !important;
        border: 1px solid rgba(255, 255, 255, 0.1) !important;
        color: white !important;
        border-radius: 12px !important;
    }
    [data-testid="stTextInput"] > div > div:focus-within, 
    [data-testid="stTextArea"] > div > div:focus-within {
        border-color: #8e44ad !important;
        box-shadow: 0 0 0 1px #8e44ad !important;
    }

    /* 6. CHAT INPUT (THE BOTTOM BAR) */
    [data-testid="stChatInput"] {
        background-color: transparent !important;
        padding-bottom: 15px;
    }
    [data-testid="stChatInput"] > div {
        background-color: rgba(255, 255, 255, 0.05) !important;
        border: 1px solid rgba(255, 255, 255, 0.1) !important;
        border-radius: 15px !important;
        color: white !important;
        padding: 2px;
    }
    
    /* === FIX 1: PUSH TEXT INSIDE INPUT TO THE RIGHT === */
    [data-testid="stChatInput"] textarea {
        padding-left: 15px !important; /* This moves the typing cursor and placeholder */
    }

    /* 7. GLASS PANELS */
    .glass-panel {
        background: rgba(255, 255, 255, 0.03);
        backdrop-filter: blur(16px);
        -webkit-backdrop-filter: blur(16px);
        border: 1px solid rgba(255, 255, 255, 0.08);
        border-radius: 20px;
        padding: 24px;
        box-shadow: 0 8px 32px rgba(0, 0, 0, 0.4);
        margin-bottom: 24px;
    }

    /* 8. CUSTOM SPINNER */
    @keyframes spin { 0% { transform: rotate(0deg); } 100% { transform: rotate(360deg); } }
    .custom-loader {
        border: 4px solid rgba(255, 255, 255, 0.1);
        border-radius: 50%;
        width: 24px; height: 24px;
        animation: spin 1s linear infinite;
        display: inline-block; vertical-align: middle; margin-right: 10px;
    }
    
    /* 9. CHAT BUBBLES */
    [data-testid="stChatMessage"] {
        background-color: rgba(255, 255, 255, 0.03);
        border: 1px solid rgba(255, 255, 255, 0.05);
        border-radius: 15px;
        margin-bottom: 12px;
    }
    
    /* === FIX: Move Bubble Text Right === */
    [data-testid="stChatMessageContent"] {
        padding-left: 20px !important;
        padding-right: 10px !important;
    }

    [data-testid="stChatMessageAvatarUser"] { background-color: #667eea; }
    [data-testid="stChatMessageAvatarAssistant"] { background-color: #764ba2; }

    header {visibility: hidden;}
    footer {visibility: hidden;}
    </style>
    """, unsafe_allow_html=True)

# --- BACKEND LOGIC ---
def extract_content(pptx_file):
    try:
        prs = Presentation(pptx_file)
        extracted_lines = []
        for s_idx, slide in enumerate(prs.slides):
            for sh_idx, shape in enumerate(slide.shapes):
                if not shape.has_text_frame: continue
                for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        extracted_lines.append(f"{{S{s_idx}:Sh{sh_idx}:P{p_idx}}} || {text}")
        return prs, "\n".join(extracted_lines)
    except Exception as e:
        return None, str(e)

def apply_changes(prs, modified_text):
    updates = {}
    for line in modified_text.split('\n'):
        # Robust regex to catch the ID and the content
        match = re.search(r"(\{S\d+:Sh\d+:P\d+\})\s*\|\|\s*(.*)", line)
        if match:
            uid = match.group(1)
            new_text = match.group(2).strip()
            parts = uid.strip("{}").split(":")
            try:
                s_idx = int(parts[0][1:])
                sh_idx = int(parts[1][2:])
                p_idx = int(parts[2][1:])
                updates[(s_idx, sh_idx, p_idx)] = new_text
            except:
                continue

    for s_idx, slide in enumerate(prs.slides):
        for sh_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame: continue
            for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                key = (s_idx, sh_idx, p_idx)
                if key in updates:
                    new_text = updates[key]
                    if len(paragraph.runs) > 0:
                        paragraph.runs[0].text = new_text
                        for i in range(1, len(paragraph.runs)):
                            paragraph.runs[i].text = ""
                    else:
                        paragraph.add_run().text = new_text
    return prs

def render_status_spinner(text, color_hex):
    return f"""
        <div style="display: flex; align-items: center; justify-content: center; padding: 20px; 
                    background: rgba(255,255,255,0.02); border-radius: 12px; border: 1px solid rgba(255,255,255,0.05); margin-top: 20px;">
            <div class="custom-loader" style="border-top: 3px solid {color_hex}; width: 30px; height: 30px; margin-right: 15px;"></div>
            <span style="font-size: 1.1rem; font-weight: 500; color: #eee; letter-spacing: 0.5px;">{text}</span>
        </div>
    """

def stream_with_initial_loader(context_prompt, placeholder_obj):
    """
    Fixed Stream Handler: Does NOT block on next() manually.
    """
    placeholder_obj.markdown(
        render_status_spinner("Connecting to Model...", "#f1c40f"), 
        unsafe_allow_html=True
    )
    
    try:
        stream = ollama.chat(
            model=MODEL_NAME,
            messages=[{'role': 'user', 'content': context_prompt}],
            stream=True
        )
        
        has_started = False
        
        for chunk in stream:
            content = chunk['message']['content']
            if not has_started:
                placeholder_obj.empty() # Clear spinner on first byte
                has_started = True
            
            yield content
            
        if not has_started:
            placeholder_obj.error("Model returned no response.")
            
    except Exception as e:
        placeholder_obj.error(f"Connection Error: {str(e)}")

# --- MAIN APP ---
st.set_page_config(page_title="Vixip Studio", layout="wide")
inject_derek_central_css()

# --- HEADER ---
st.markdown("""
<div class="glass-panel" style="text-align: center; padding: 30px;">
    <h1 class="aurora-text">Vixip Studio</h1>
    <p style="opacity: 0.6; margin-top: 5px; font-size: 0.9rem;">AI Presentation Architect</p>
</div>
""", unsafe_allow_html=True)

# --- SESSION STATE ---
if "chat_history" not in st.session_state: st.session_state["chat_history"] = []
if "raw_pptx_text" not in st.session_state: st.session_state["raw_pptx_text"] = ""
if "generator_instruction" not in st.session_state: st.session_state["generator_instruction"] = ""
if "file_uploaded" not in st.session_state: st.session_state["file_uploaded"] = False

# --- UPLOAD SECTION ---
if not st.session_state["file_uploaded"]:
    st.markdown("### Upload Slides (.pptx)")
    # === FIX: Added non-empty label + label_visibility to fix Console Warning ===
    uploaded_file = st.file_uploader("Upload PPTX", type="pptx", label_visibility="collapsed")
    
    if uploaded_file:
        with open("temp.pptx", "wb") as f:
            f.write(uploaded_file.getbuffer())
        
        prs, text = extract_content("temp.pptx")
        if prs:
            st.session_state["raw_pptx_text"] = text
            st.session_state["file_uploaded"] = True
            st.rerun()
else:
    with st.sidebar:
        st.success("Analysis Active")
        if st.button("Upload New File"):
            st.session_state["file_uploaded"] = False
            st.session_state["raw_pptx_text"] = ""
            st.rerun()

# --- MAIN INTERFACE ---
if st.session_state["file_uploaded"]:
    
    tab_chat, tab_gen = st.tabs(["Chat & Strategy", "Slide Generator"])

    # === TAB 1: STREAMING CHAT ===
    with tab_chat:
        chat_container = st.container()
        
        with chat_container:
            if not st.session_state["chat_history"]:
                st.info("👋 Hey! I've loaded your slideshow. Wanna discuss your slides?")
            
            for msg in st.session_state["chat_history"]:
                with st.chat_message(msg["role"]):
                    st.write(msg["content"])

        if prompt := st.chat_input("Ask about your slides..."):
            st.session_state["chat_history"].append({"role": "user", "content": prompt})
            with chat_container:
                with st.chat_message("user"):
                    st.write(prompt)

            with chat_container:
                with st.chat_message("assistant"):
                    stream_placeholder = st.empty()
                    context = f"PPTX Content:\n{st.session_state['raw_pptx_text']}\n\nUser Question: {prompt}"
                    full_response = st.write_stream(stream_with_initial_loader(context, stream_placeholder))
                    st.session_state["chat_history"].append({"role": "assistant", "content": full_response})

        if st.session_state["chat_history"]:
            st.markdown("---")
            if st.button("✨ Use this strategy in Generator"):
                last_msg = st.session_state["chat_history"][-1]["content"]
                st.session_state["generator_instruction"] = (
                    "Based on our previous chat, implement these changes.\n"
                    f"CONTEXT: {last_msg}"
                )
                st.toast("Strategy copied!", icon="🚀")

    # === TAB 2: GENERATOR (With Fail-Safe) ===
    with tab_gen:
        instruction = st.text_area(
            "Instructions", 
            value=st.session_state["generator_instruction"],
            placeholder="e.g. Translate to Korean and make the tone more professional...",
            height=150
        )
        
        col1, col2 = st.columns([1, 4])
        with col1:
            run_btn = st.button("Run Transformation", type="primary", use_container_width=True)
        
        if run_btn:
            if not instruction:
                st.warning("Please enter instructions.")
            else:
                status_box = st.empty()
                
                # STRICTER PROMPT
                sys_prompt = (
                    "You are a Professional Presentation Architect. \n"
                    "NOTE: Please remain grounded, impartial, realistic and accurate and remember the user does not seet the internal pptx positional information. Only you do. Keep track of any connstraitns, requirements, preferences, specifications, etc forever. \n"
                    "1. PLAN: First, think, reason and plan the changes. Consider what the user wants, the best way to achieve it along with the best way to get there. Also consider what is already in the slide. Think, reason, plan, draft and consider multiple possibilities given the prior then tackle the whole thing.\n"
                    f"2. SEPARATOR: When ready to output code, you MUST output this line: {SEPARATOR} ENSURE YOU USE {SEPARATOR} otherwise the content might not be parsed!\n"
                    "3. OUTPUT: After that line, write the modified lines in format {S#:Sh#:P#} || Text\n"
                    "Do not change IDs and when discussing slides do not mention them since the user does not see the {S1:Sh2:P0} || they only see the stuff after the ||. Only change the text not anything before the ||."
                )
                full_prompt = f"INSTRUCTION: {instruction}\n\nDATA:\n{st.session_state['raw_pptx_text']}"
                
                try:
                    status_box.markdown(render_status_spinner("Initializing Agent...", "#f1c40f"), unsafe_allow_html=True)
                    
                    stream = ollama.chat(
                        model=MODEL_NAME,
                        messages=[{'role': 'system', 'content': sys_prompt}, {'role': 'user', 'content': full_prompt}],
                        stream=True
                    )
                    
                    status_box.markdown(render_status_spinner("Working... (this might take awhile...)", "#8e44ad"), unsafe_allow_html=True)
                    
                    thinking_buffer = ""
                    slide_content = ""
                    mode = "THINKING"
                    
                    for chunk in stream:
                        content = chunk['message']['content']
                        
                        if mode == "THINKING":
                            thinking_buffer += content
                            # Check buffer for separator
                            if SEPARATOR in thinking_buffer:
                                mode = "WRITING"
                                status_box.markdown(render_status_spinner("Generating Final Slides...", "#2ecc71"), unsafe_allow_html=True)
                                
                                # Process what came after separator in this specific chunk
                                parts = thinking_buffer.split(SEPARATOR)
                                if len(parts) > 1:
                                    slide_content += parts[1]
                        else:
                            slide_content += content

                    # === FAIL-SAFE: If model finished but never output separator ===
                    if mode == "THINKING":
                        # Recover slide text if it exists in the thinking buffer
                        if "{S0" in thinking_buffer or "||" in thinking_buffer:
                            # Try to strip out the conversational part roughly
                            lines = thinking_buffer.split('\n')
                            valid_lines = [l for l in lines if "||" in l and "{" in l]
                            if valid_lines:
                                slide_content = "\n".join(valid_lines)
                                st.warning("Model ignored the separator, but slide content was recovered.")
                            else:
                                st.error("Model finished but produced no valid slide lines.")
                        else:
                            st.error("Model finished but did not generate any slide updates.")
                    
                    # FINAL PROCESSING
                    if slide_content.strip():
                        prs = Presentation("temp.pptx")
                        updated = apply_changes(prs, slide_content)
                        out = BytesIO()
                        updated.save(out)
                        out.seek(0)
                        
                        status_box.empty()
                        st.success("✅ Transformation Complete!")
                        st.download_button("📥 Download Enhanced PPTX", out, "enhanced.pptx", type="primary")
                    
                except Exception as e:
                    status_box.error(f"Error: {e}")
